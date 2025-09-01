import os
import json
import fitz  # PyMuPDF
import traceback
from typing import List, Dict, Any, Optional
from pathlib import Path
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.messages import HumanMessage
from pptx import Presentation
from prompts import pdf_analysis_prompt,email_analysis_prompt,call_analysis_prompt
# import aspose.slides as slides
import tempfile 
# Setup API key
os.environ["GOOGLE_API_KEY"] = "your_google_api_key"

class StartupAnalyzer:
    def __init__(self):
        # Initialize text-only model
        self.text_model = ChatGoogleGenerativeAI(
            model="gemini-1.5-flash",
            temperature=0.1
        )
    # def convert_pptx_to_pdf(self, pptx_path: str) -> str:

    #     """Convert PPTX to PDF using Aspose.Slides."""
    #     try:
            
            
    #         print(f"🔄 Converting {pptx_path} to PDF using Aspose...")
            
    #         # Load presentation
    #         presentation = slides.Presentation(pptx_path)
            
    #         # Create temporary PDF path
    #         temp_pdf = tempfile.NamedTemporaryFile(suffix='.pdf', delete=False)
    #         temp_pdf.close()
            
    #         # Save as PDF
    #         presentation.save(temp_pdf.name, slides.export.SaveFormat.PDF)
            
    #         print(f"✅ Successfully converted PPTX to PDF")
    #         return temp_pdf.name
            
    #     except ImportError:
    #         raise Exception("aspose-slides not installed. Run: pip install aspose-slides")
    #     except Exception as e:
    #         raise Exception(f"Aspose conversion failed: {e}")


    def extract_text_from_pdf(self, pdf_path: str) -> str:
        """Extract text from PDF file."""
        try:
            doc = fitz.open(pdf_path)
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text
        except Exception as e:
            return f"Error extracting PDF: {str(e)}"

    def create_pdf_analysis_prompt(self, pdf_text: str) -> str:
        """Create detailed information extraction prompt for PDF documents."""
        return pdf_analysis_prompt.format(pdf_text=pdf_text)

    def create_email_analysis_prompt(self, raw_email_text: str) -> str:
        """Create detailed information extraction prompt for raw email text."""
        return email_analysis_prompt.format(raw_email_text=raw_email_text)


    def create_call_analysis_prompt(self, raw_transcript: str) -> str:
        """Create detailed information extraction prompt for raw call transcript."""
        return call_analysis_prompt.format(raw_transcript=raw_transcript)

    def analyze_pdf_document(self, pdf_path: str, doc_type: str = "general") -> Dict[str, Any]:
        """Analyze PDF document (pitch deck, financials, legal docs)."""
        try:
            # Extract text
            pdf_text = self.extract_text_from_pdf(pdf_path)
            if pdf_text.startswith("Error"):
                return {"error": pdf_text, "status": "failed"}

            # Create prompt
            prompt = self.create_pdf_analysis_prompt(pdf_text)
            
            # Get analysis
            messages = [HumanMessage(content=prompt)]
            response = self.text_model.invoke(messages)
            
            return {
                "document_type": doc_type,
                "file_path": pdf_path,
                "analysis": response.content,
                "status": "success"
            }
            
        except Exception as e:
            return {
                "error": f"Analysis failed: {str(e)}",
                "traceback": traceback.format_exc(),
                "status": "failed"
            }

    # def analyze_pitch_deck_pptx(self, pptx_path: str) -> Dict[str, Any]:
    #     """Analyze PowerPoint by converting to PDF first using Aspose."""
    #     temp_pdf_path = None
        
    #     try:
    #         print(f"\n🔄 Processing PowerPoint file: {pptx_path}")
            
    #         # Step 1: Convert PPTX to PDF using Aspose
    #         print("📄 Converting PPTX to PDF...")
    #         temp_pdf_path = self.convert_pptx_to_pdf(pptx_path)
            
    #         # Step 2: Extract text from PDF using fitz
    #         print("📝 Extracting text from PDF...")
    #         extracted_text = self.extract_text_from_pdf(temp_pdf_path)
            
    #         if extracted_text.startswith("Error"):
    #             return {
    #                 "document_type": "pitch_deck_pptx",
    #                 "file_path": pptx_path,
    #                 "error": extracted_text,
    #                 "status": "failed"
    #             }
            
    #         print(f"✅ Extracted {len(extracted_text)} characters of text")
    #         print(f"📝 Text preview (first 300 chars):\n{extracted_text[:300]}...")
            
    #         # Step 3: Analyze with AI
    #         if len(extracted_text.strip()) < 100:
    #             return {
    #                 "document_type": "pitch_deck_pptx",
    #                 "file_path": pptx_path,
    #                 "error": f"Insufficient text content extracted ({len(extracted_text)} chars)",
    #                 "raw_text": extracted_text,
    #                 "status": "failed"
    #             }
            
    #         prompt = self.create_pdf_extraction_prompt(extracted_text)
            
    #         print("🤖 Sending to AI model for analysis...")
            
    #         messages = [HumanMessage(content=prompt)]
    #         response = self.text_model.invoke(messages)
            
    #         return {
    #             "document_type": "pitch_deck_pptx",
    #             "file_path": pptx_path,
    #             "extracted_information": response.content,
    #             "raw_text": extracted_text,
    #             "conversion_method": "aspose_pptx_to_pdf",
    #             "status": "success"
    #         }
            
    #     except Exception as e:
    #         return {
    #             "document_type": "pitch_deck_pptx",
    #             "file_path": pptx_path,
    #             "error": f"PPTX processing failed: {str(e)}",
    #             "traceback": traceback.format_exc(),
    #             "status": "failed"
    #         }
    

    def analyze_raw_email(self, raw_email_text: str) -> Dict[str, Any]:
        """Analyze raw email text."""
        try:
            prompt = self.create_email_analysis_prompt(raw_email_text)
            
            messages = [HumanMessage(content=prompt)]
            response = self.text_model.invoke(messages)
            
            return {
                "document_type": "email",
                "analysis": response.content,
                "status": "success"
            }
            
        except Exception as e:
            return {
                "error": f"Email analysis failed: {str(e)}",
                "traceback": traceback.format_exc(),
                "status": "failed"
            }

    def analyze_raw_call_transcript(self, raw_transcript: str) -> Dict[str, Any]:
        """Analyze raw call transcript text."""
        try:
            prompt = self.create_call_analysis_prompt(raw_transcript)
            
            messages = [HumanMessage(content=prompt)]
            response = self.text_model.invoke(messages)
            
            return {
                "document_type": "call_transcript",
                "analysis": response.content,
                "status": "success"
            }
            
        except Exception as e:
            return {
                "error": f"Call analysis failed: {str(e)}",
                "traceback": traceback.format_exc(),
                "status": "failed"
            }

    # def extract_pptx_text(self, pptx_path: str) -> str:
    #     """Extract all text from PowerPoint slides with better error handling."""
    #     try:
    #         print(f"Attempting to extract text from: {pptx_path}")
            
    #         if not os.path.exists(pptx_path):
    #             return f"Error: File not found at {pptx_path}"
            
    #         prs = Presentation(pptx_path)
    #         all_text = ""
    #         slide_count = len(prs.slides)
            
    #         print(f"Found {slide_count} slides in presentation")
            
    #         if slide_count == 0:
    #             return "Error: No slides found in presentation"
            
    #         for i, slide in enumerate(prs.slides):
    #             slide_text = f"\n--- SLIDE {i + 1} ---\n"
    #             shape_count = 0
                
    #             # Extract text from all shapes
    #             for shape in slide.shapes:
    #                 shape_count += 1
    #                 try:
    #                     if hasattr(shape, "text") and shape.text.strip():
    #                         slide_text += f"Shape {shape_count}: {shape.text.strip()}\n"
    #                     elif hasattr(shape, "text_frame") and shape.text_frame:
    #                         for paragraph in shape.text_frame.paragraphs:
    #                             para_text = ""
    #                             for run in paragraph.runs:
    #                                 para_text += run.text
    #                             if para_text.strip():
    #                                 slide_text += f"Shape {shape_count}: {para_text.strip()}\n"
    #                 except Exception as shape_error:
    #                     print(f"Error extracting text from shape {shape_count} on slide {i+1}: {shape_error}")
    #                     continue
                
    #             if slide_text.strip() != f"--- SLIDE {i + 1} ---":
    #                 all_text += slide_text
    #             else:
    #                 all_text += f"\n--- SLIDE {i + 1} ---\n[No readable text found on this slide]\n"
            
    #         print(f"Extracted text length: {len(all_text)} characters")
            
    #         if len(all_text.strip()) < 50:  # Very little content extracted
    #             return f"Warning: Only minimal text extracted from presentation. Extracted content:\n{all_text}"
            
    #         return all_text
            
    #     except Exception as e:
    #         error_msg = f"Error extracting PPTX text: {str(e)}"
    #         print(error_msg)
    #         return error_msg
# Simplified usage functions
def analyze_startup_documents(file_paths: List[str]) -> Dict[str, Any]:
    """Analyze multiple startup documents."""
    analyzer = StartupAnalyzer()
    results = {}
    
    for file_path in file_paths:
        file_path = Path(file_path)
        
        try:
            if file_path.suffix.lower() == '.pdf':
                result = analyzer.analyze_pdf_document(str(file_path))
            elif file_path.suffix.lower() in ['.pptx', '.ppt']:
                result = analyzer.analyze_pitch_deck_pptx(str(file_path))
            else:
                result = {"error": f"Unsupported file type: {file_path.suffix}"}
            
            results[file_path.name] = result
            
        except Exception as e:
            results[file_path.name] = {
                "error": f"Processing failed: {str(e)}",
                "traceback": traceback.format_exc()
            }
    
    return results

def analyze_raw_email_text(raw_email: str) -> Dict[str, Any]:
    """Analyze raw email text."""
    analyzer = StartupAnalyzer()
    return analyzer.analyze_raw_email(raw_email)

def analyze_raw_call_text(raw_transcript: str) -> Dict[str, Any]:
    """Analyze raw call transcript text."""
    analyzer = StartupAnalyzer()
    return analyzer.analyze_raw_call_transcript(raw_transcript)

# Example usage
if __name__ == "__main__":
    # Example 1: Analyze startup documents
    documents = [
        "startup_deck.pdf",
    ]
    
    results = analyze_startup_documents(documents)
    
    for filename, analysis in results.items():
        print(f"\n{'='*50}")
        print(f"📄 {filename}")
        print(f"{'='*50}")
        
        if analysis.get('status') == 'success':
            print(analysis['analysis'])
        else:
            print(f"❌ Error: {analysis.get('error')}")
    
    # Example 2: Analyze raw email
    raw_email = """
From: john.founder@startup.com
To: sarah.investor@vcfund.com
Subject: Monthly Updates & Follow-ups
Date: March 15, 2024, 2:30 PM

========== EMAIL 1 ==========
From: john.founder@startup.com
To: sarah.investor@vcfund.com
Subject: February Update - TechStartup Inc.
Date: February 28, 2024, 4:15 PM

Hi Sarah,

February update for TechStartup Inc:

FINANCIALS:
- Revenue: $38K MRR (up 12% from January)
- Burn rate: $45K/month
- Runway: 16 months remaining
- New customers: 23 this month

PRODUCT:
- Launched mobile app (iOS and Android)
- Integration with Salesforce completed
- 89% customer satisfaction score

TEAM:
- Hired senior engineer from Google
- Team size now 15 people
- Opening sales office in Austin

CHALLENGES:
- Competition from BigTech Corp increasing
- Customer acquisition costs rising

Best,
John

========== EMAIL 2 ==========
From: sarah.investor@vcfund.com
To: john.founder@startup.com
Subject: RE: February Update - Questions
Date: March 2, 2024, 10:22 AM

John,

Thanks for the update. Few questions:

1. What's driving the CAC increase? Can you share specific numbers?
2. How are you planning to compete with BigTech Corp?
3. Can you send Q1 financial projections?
4. What's the timeline for Series A fundraising?

Also, our investment committee wants to schedule a deep-dive call for next week.

Best,
Sarah

========== EMAIL 3 ==========
From: john.founder@startup.com
To: sarah.investor@vcfund.com
Subject: RE: February Update - Answers & March Update
Date: March 15, 2024, 2:30 PM

Sarah,

Answering your questions + March mid-month update:

ANSWERS TO YOUR QUESTIONS:
1. CAC Analysis:
   - January CAC: $145
   - February CAC: $189 (+30%)
   - Main driver: Google Ads cost increased 40%, Facebook ads 25%
   - New strategy: Focusing on content marketing and referrals

2. BigTech Competition Strategy:
   - They're targeting enterprise ($50K+ deals)
   - We're doubling down on SMB market ($500-5K deals)
   - Our advantage: 10x faster implementation time
   - Building deeper integrations they can't match

3. Q1 Projections (attached spreadsheet):
   - March revenue target: $42K MRR
   - Q1 total revenue: $125K
   - Q1 ending MRR: $45K

4. Series A Timeline:
   - Planning to start in June 2024
   - Target raise: $8-12M
   - Want to hit $60K MRR before starting

MARCH UPDATE (First half):
- Revenue: $41K MRR (hit target early!)
- New customers: 31 (best month ever)
- Churn rate: Only 2.1% (down from 4.3%)
- Enterprise pilot with Fortune 500 company started

TEAM NEWS:
- VP of Sales starts Monday (ex-Salesforce)
- Engineering team now 8 people
- Opened Austin office officially

DEEP-DIVE CALL:
Next week works great. How about Wednesday 2 PM PST?

Best,
John

P.S. - Just closed our biggest deal ever: $15K/year contract!

========== EMAIL 4 ==========
From: sarah.investor@vcfund.com
To: john.founder@startup.com
Subject: RE: Great progress! Next steps
Date: March 18, 2024, 9:45 AM

John,

Excellent progress! The numbers look strong and I like the competition strategy.

Wednesday 2 PM works. I'll send calendar invite.

For the call, please prepare:
- Customer cohort analysis
- Detailed CAC/LTV breakdown by channel
- Competitive analysis deep-dive
- Series A planning timeline
- Use of funds projection

Also bringing our technical partner to evaluate your product architecture.

Talk Wednesday,
Sarah

========== EMAIL 5 ==========
From: john.founder@startup.com
To: sarah.investor@vcfund.com
Subject: Post-call follow-up & April preview
Date: March 22, 2024, 6:30 PM

Sarah,

Great call Wednesday! Attaching the additional materials you requested:

ATTACHMENTS:
- Cohort_Analysis_Q1_2024.xlsx
- Technical_Architecture_Overview.pdf
- Competitive_Landscape_Analysis.pptx

APRIL PREVIEW:
Early indicators show we might hit $47K MRR by month-end (ahead of plan).

Two big updates:
1. Fortune 500 pilot going extremely well - they want to discuss enterprise contract
2. Received acquisition inquiry from MidSize Corp ($45M offer - not interested, but validates market)

SERIES A UPDATE:
Based on our momentum, considering moving Series A timeline up to May. Thoughts?

Best,
John
"""

    # email_result = analyze_raw_email_text(raw_email)
    
    # print(f"\n{'='*50}")
    # print("📧 EMAIL ANALYSIS")
    # print(f"{'='*50}")
    # print(email_result.get('analysis', email_result.get('error')))
    
    # Example 3: Analyze raw call transcript
    raw_call = """
INVESTOR-FOUNDER CALL SESSION
Date: March 20, 2024
Total Duration: 2 hours 15 minutes
Participants: John Founder (CEO), Mike Cofounder (CTO), Sarah Investor (Partner), David Tech (Technical DD)

========== CALL 1: BUSINESS OVERVIEW (45 minutes) ==========
Time: 10:00 AM - 10:45 AM
Participants: John Founder, Sarah Investor

Sarah: John, thanks for making time. Let's start with your current business metrics.

John: Absolutely. We're at $41K MRR as of mid-March, growing about 18% month-over-month. We have 847 active users, with 312 paying customers.

Sarah: What's your unit economics looking like?

John: CAC is averaging $165 across all channels, though it varies significantly. Google Ads runs about $210, but referrals are only $45. LTV is $2,400 based on 24-month average customer lifecycle.

Sarah: That's roughly 14:1 LTV to CAC ratio. How confident are you in those LTV calculations?

John: Pretty confident. We've got 18 months of cohort data now. Our oldest cohorts actually show higher LTV - around $2,800 - so we might be conservative.

Sarah: Tell me about your competitive landscape.

John: There are three main competitors. BigCorp has 60% market share but focuses on enterprise. StartupX raised $50M last year but their product is more complex. LocalPlayer is our closest competitor but they're only in California.

Sarah: How do you differentiate?

John: Speed and simplicity. Our implementation takes 2 hours versus their 2 weeks. Our pricing is transparent - $99/month vs their complex enterprise pricing.

Sarah: What about defensibility?

John: We're building network effects through our marketplace feature. Plus we have proprietary algorithms that took 2 years to develop. Filing patents on our core IP.

Sarah: Walk me through your financial projections.

John: We're targeting $60K MRR by June, $100K by December. Revenue growth should accelerate as we launch enterprise tier and marketplace features.

Sarah: What's your burn rate and runway?

John: Currently burning $52K monthly. We have $850K in the bank, so about 16 months runway. But revenue is growing faster than burn, so effective runway is longer.

========== CALL 2: TECHNICAL DEEP DIVE (50 minutes) ==========
Time: 11:00 AM - 11:50 AM
Participants: Mike Cofounder (CTO), David Tech (Technical Partner)

David: Mike, I'd like to understand your technical architecture.

Mike: Sure. We're built on AWS using microservices architecture. React frontend, Node.js backend, PostgreSQL database with Redis for caching.

David: How do you handle scalability?

Mike: We use auto-scaling groups for compute, RDS with read replicas for database scaling. We've load-tested up to 50,000 concurrent users without issues.

David: What about data security?

Mike: SOC 2 Type II compliant, all data encrypted at rest and in transit. We use AWS KMS for key management and have regular penetration testing.

David: Tell me about your development process.

Mike: Agile methodology, 2-week sprints. Full CI/CD pipeline with automated testing. Code coverage is 89%, aiming for 95%.

David: How's your team structured?

Mike: 8 engineers total: 3 frontend, 3 backend, 1 DevOps, 1 mobile. Planning to hire 2 more senior engineers next quarter.

David: What technical challenges are you facing?

Mike: Main challenge is handling real-time data processing as we scale. We're migrating to event-driven architecture using AWS EventBridge.

David: How do you handle customer data integration?

Mike: We have pre-built connectors for 15 major platforms. Custom integrations typically take 1-2 weeks to build. API-first architecture makes it straightforward.

David: What's your uptime and performance metrics?

Mike: 99.97% uptime over last 12 months. Average API response time is 120ms. We have comprehensive monitoring with PagerDuty alerts.

========== CALL 3: MARKET & GROWTH STRATEGY (40 minutes) ==========
Time: 1:00 PM - 1:40 PM
Participants: John Founder, Sarah Investor

Sarah: Let's talk about your go-to-market strategy.

John: We have three main channels. Inbound marketing generates 40% of leads through content and SEO. Outbound sales handles 35%, mostly enterprise prospects. Referrals are 25% and our highest-converting channel.

Sarah: What's your sales process?

John: SMB deals: 7-day average sales cycle, mostly self-serve with email support. Enterprise: 45-day cycle with demos and pilot programs. Conversion rates are 23% SMB, 12% enterprise.

Sarah: How big is your addressable market?

John: TAM is $12 billion globally. SAM focusing on US SMBs is $2.8 billion. Our SOM over next 5 years is realistically $280 million.

Sarah: Those are big numbers. How did you calculate them?

John: Used bottom-up analysis. 2.3 million SMBs in our target segments, average spend of $1,200 annually on solutions like ours. Cross-referenced with Gartner and McKinsey reports.

Sarah: What's your customer acquisition strategy?

John: Focusing on product-led growth. Free tier drives adoption, then convert to paid. We're also building partnership channel - already have 12 resellers signed up.

Sarah: Tell me about customer feedback.

John: NPS score is 67. Common feedback: easy to use, great support, saves significant time. Main complaints: want more advanced reporting and mobile app improvements.

Sarah: How do you plan to expand?

John: Three phases: First, geographic expansion to Canada and UK. Second, adjacent products for existing customers. Third, move upmarket to mid-enterprise.

========== CALL 4: FINANCIALS & FUNDING (20 minutes) ==========
Time: 1:45 PM - 2:05 PM
Participants: John Founder, Sarah Investor

Sarah: Let's discuss your Series A plans.

John: Planning to raise $8-12 million in May-June timeframe. Want to hit $50K MRR first to improve valuation.

Sarah: What's the use of funds?

John: 50% engineering team expansion, 30% sales and marketing, 15% operations and infrastructure, 5% working capital buffer.

Sarah: What valuation range are you targeting?

John: Based on comparable companies and our metrics, thinking $35-45 million pre-money. Open to feedback on that.

Sarah: How much runway will that give you?

John: Should get us to 36+ months runway and profitability. Planning to be cash-flow positive by month 18-20 of the funding.

Sarah: Any other investors you're talking to?

John: Early conversations with GrowthVC and TechFund. Haven't started formal process yet. You're our preferred partner based on portfolio fit and expertise.

Sarah: What milestones do you need before fundraising?

John: Hit $50K MRR, complete enterprise product tier, sign 2-3 marquee enterprise customers, expand team to 25 people.

========== WRAP-UP DISCUSSION (10 minutes) ==========
Time: 2:05 PM - 2:15 PM
All participants

Sarah: This has been really comprehensive. David, your initial technical assessment?

David: Strong technical foundation, experienced team, good architecture decisions. Some scalability work needed but nothing concerning.

Sarah: John, what are your immediate next steps?

John: Focus on hitting $50K MRR by end of April, close the Fortune 500 pilot deal, and finalize our enterprise pricing strategy.

Sarah: Perfect. We'll do internal review over the next week. I'll have feedback for you by early next week, and if we decide to move forward, we can discuss term sheets.

John: That sounds great. Really appreciate the thorough process and everyone's time today.

Mike: Thanks David for the technical deep dive. Helpful to get external validation of our approach.

Sarah: We'll be in touch soon. Great meeting everyone.

[END OF CALL SESSION]
"""

    
    # call_result = analyze_raw_call_text(raw_call)
    
    # print(f"\n{'='*50}")
    # print("📞 CALL ANALYSIS") 
    # print(f"{'='*50}")
    # print(call_result.get('analysis', call_result.get('error')))

